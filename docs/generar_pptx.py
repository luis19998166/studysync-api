"""
Genera la presentación PowerPoint de StudySync para la defensa final.
Tema oscuro, espacio para capturas, flujo completo en vivo.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.enum.dml import MSO_THEME_COLOR
import copy

OUT = r"d:\repos\studysync-api\docs\StudySync_Presentacion_Final.pptx"

# ── Paleta ─────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0F, 0x17, 0x2A)   # azul muy oscuro
SURFACE  = RGBColor(0x1E, 0x29, 0x3B)   # azul oscuro (tarjetas)
PRIMARY  = RGBColor(0x63, 0x66, 0xF1)   # violeta (acento)
SUCCESS  = RGBColor(0x22, 0xC5, 0x5E)   # verde
WARNING  = RGBColor(0xF5, 0x9E, 0x0B)   # ámbar
DANGER   = RGBColor(0xEF, 0x44, 0x44)   # rojo
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
MUTED    = RGBColor(0x94, 0xA3, 0xB8)
YELLOW   = RGBColor(0xFA, 0xCC, 0x15)
CYAN     = RGBColor(0x06, 0xB6, 0xD4)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # completamente en blanco

def slide():
    s = prs.slides.add_slide(BLANK)
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return s

def box(s, x, y, w, h, bg=None, alpha=None):
    shape = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    if bg:
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg
    else:
        shape.fill.background()
    return shape

def txt(s, text, x, y, w, h,
        size=24, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return tb

def header_bar(s, title, subtitle="", pill_text="", pill_color=PRIMARY):
    """Barra superior con título y pastilla de estado."""
    box(s, 0, 0, 13.33, 1.1, bg=SURFACE)
    box(s, 0, 1.08, 13.33, 0.04, bg=PRIMARY)
    txt(s, title, 0.35, 0.1, 9, 0.8, size=28, bold=True, color=WHITE)
    if subtitle:
        txt(s, subtitle, 0.35, 0.65, 9, 0.5, size=13, color=MUTED)
    if pill_text:
        b = box(s, 10.8, 0.25, 2.2, 0.55, bg=pill_color)
        b.line.fill.background()
        txt(s, pill_text, 10.8, 0.25, 2.2, 0.55, size=12, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)

def screenshot_placeholder(s, x, y, w, h, label="[ Captura de pantalla ]"):
    """Rectángulo punteado para pegar captura."""
    ph = box(s, x, y, w, h, bg=RGBColor(0x1A, 0x28, 0x3A))
    ph.line.color.rgb = PRIMARY
    ph.line.width = Pt(1.5)
    txt(s, label, x, y + h/2 - 0.25, w, 0.5,
        size=13, color=PRIMARY, align=PP_ALIGN.CENTER, italic=True)
    return ph

def num_badge(s, num, x, y, color=PRIMARY):
    """Círculo numerado."""
    c = box(s, x, y, 0.45, 0.45, bg=color)
    txt(s, str(num), x, y, 0.45, 0.45, size=14, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)

def bullet_row(s, icon, text, x, y, w=12, icon_color=SUCCESS, size=15):
    txt(s, icon, x, y, 0.4, 0.4, size=size, color=icon_color, align=PP_ALIGN.CENTER)
    txt(s, text, x+0.42, y, w-0.42, 0.4, size=size, color=WHITE)

def footer(s, num, total):
    txt(s, f"StudySync API  •  Programación IV — UPDS 2026  •  {num}/{total}",
        0, 7.15, 13.33, 0.35, size=10, color=MUTED, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 1 — PORTADA
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
box(s, 0, 0, 13.33, 7.5, bg=BG)
# Degradado visual con cajas
box(s, 0, 0, 13.33, 3.8, bg=SURFACE)
box(s, 0, 3.78, 13.33, 0.06, bg=PRIMARY)

txt(s, "StudySync API", 0.5, 0.5, 12.3, 1.6,
    size=58, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Proyecto Final — Programación IV", 0.5, 2.0, 12.3, 0.8,
    size=24, color=PRIMARY, align=PP_ALIGN.CENTER)
txt(s, "Verificación completa del flujo en producción", 0.5, 2.65, 12.3, 0.7,
    size=17, color=MUTED, align=PP_ALIGN.CENTER)

# Pills de tecnología
pills = [
  ("Node.js", SUCCESS), ("Prisma + Supabase", CYAN), ("Redis Pub/Sub", WARNING),
  ("Socket.io", PRIMARY), ("JWT + Helmet", DANGER), ("CloudFormation", YELLOW)
]
px = 0.6
for label, color in pills:
    b = box(s, px, 4.15, 1.9, 0.45, bg=color)
    b.line.fill.background()
    txt(s, label, px, 4.15, 1.9, 0.45, size=12, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    px += 2.05

txt(s, "M.Sc. Jimmy Nataniel Requena Llorentty  •  UPDS Bolivia  •  2026",
    0.5, 4.85, 12.3, 0.5, size=14, color=MUTED, align=PP_ALIGN.CENTER)
txt(s, "github.com/bolivianotech/studysync-api",
    0.5, 5.35, 12.3, 0.5, size=13, color=PRIMARY, align=PP_ALIGN.CENTER)
txt(s, "studysync-api-a5o9.onrender.com",
    0.5, 5.75, 12.3, 0.5, size=13, color=CYAN, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 2 — AGENDA
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header_bar(s, "Agenda de la Defensa", "15 minutos de demostración en vivo", "15 min", SUCCESS)

items = [
  ("01", "¿Qué construimos?",         "Stack completo y arquitectura del sistema",                    PRIMARY),
  ("02", "API en producción",          "Swagger UI — registro, login y BearerAuth en vivo",            CYAN),
  ("03", "Datos reales en Supabase",   "Table Editor — ver filas creadas en tiempo real",              SUCCESS),
  ("04", "Frontend + WebSockets",      "Dos pestañas — evento Redis llega a todos los clientes",       WARNING),
  ("05", "3 Grupos de estudio",        "9+ sesiones de Ana Torres, Carlos Méndez y Sofía Vargas",     YELLOW),
  ("06", "IaC con LocalStack",         "S3 + DynamoDB + SSM creados desde un YAML",                   DANGER),
  ("07", "Repositorio GitHub",         "Commits, estructura, README y guía de instalación",            MUTED),
]

for i, (num, title, desc, color) in enumerate(items):
    y = 1.3 + i * 0.83
    box(s, 0.4, y, 12.5, 0.72, bg=SURFACE)
    b2 = box(s, 0.4, y, 0.12, 0.72, bg=color)
    b2.line.fill.background()
    txt(s, num, 0.6, y+0.12, 0.7, 0.5, size=15, bold=True, color=color)
    txt(s, title, 1.4, y+0.06, 5.5, 0.35, size=15, bold=True, color=WHITE)
    txt(s, desc,  1.4, y+0.38, 10.8, 0.28, size=11, color=MUTED)

footer(s, 2, 18)

# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 3 — STACK TECNOLÓGICO
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header_bar(s, "Stack Tecnológico", "Cada pieza tiene un propósito específico", "12 tecnologías")

tecnologias = [
  ("Node.js + Express v5",   "Servidor HTTP y rutas REST",            PRIMARY),
  ("Prisma ORM v6",          "Acceso type-safe a PostgreSQL",          CYAN),
  ("Supabase",               "PostgreSQL en la nube (gratis)",         SUCCESS),
  ("Redis / Upstash",        "Pub/Sub para mensajes en tiempo real",   WARNING),
  ("Socket.io v4",           "WebSockets hacia el navegador",          YELLOW),
  ("JWT + bcryptjs",         "Autenticación sin estado",               DANGER),
  ("Helmet",                 "12 headers de seguridad HTTP",           RGBColor(0xA8,0x55,0xF7)),
  ("Swagger / OpenAPI 3.0",  "Documentación interactiva de la API",    CYAN),
  ("CloudFormation",         "Infraestructura como Código (IaC)",      WARNING),
  ("LocalStack",             "Simula AWS localmente con Docker",        SUCCESS),
  ("Render.com",             "Deploy automático desde GitHub",          PRIMARY),
  ("Vanilla JS + CSS",       "Frontend sin frameworks, en public/",    MUTED),
]

for i, (tech, desc, color) in enumerate(tecnologias):
    col = i % 3
    row = i // 3
    x = 0.3 + col * 4.35
    y = 1.25 + row * 1.45
    box(s, x, y, 4.1, 1.28, bg=SURFACE)
    b = box(s, x, y, 0.1, 1.28, bg=color)
    b.line.fill.background()
    txt(s, tech, x+0.25, y+0.12, 3.7, 0.42, size=13, bold=True, color=WHITE)
    txt(s, desc, x+0.25, y+0.52, 3.7, 0.55, size=11, color=MUTED)

footer(s, 3, 18)

# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 4 — ARQUITECTURA
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header_bar(s, "Arquitectura del Sistema", "Flujo completo de una petición")

arch_text = """\
  NAVEGADOR (public/index.html)
       │  HTTP REST + WebSocket (Socket.io)
       ▼
  ┌─────────────────────────────────────────┐
  │         EXPRESS API  (src/)             │
  │  Helmet → CORS → RateLimit → JWT Auth   │
  │  POST /auth/register   POST /auth/login │
  │  GET /api/sesiones  (público)           │
  │  POST/PUT/DELETE /api/sesiones (JWT)    │
  └────────────┬────────────────┬───────────┘
               │                │
  ┌────────────▼───┐   ┌────────▼──────────┐
  │  Upstash Redis │   │  Supabase (Cloud) │
  │  pub.publish() │   │  PostgreSQL       │
  │  sub.subscribe │   │  Prisma ORM       │
  └────────────┬───┘   └───────────────────┘
               │
  ┌────────────▼──────────────┐
  │  subscribers/             │
  │  notificaciones.js        │
  │  io.emit('nuevo-evento')  │
  └────────────┬──────────────┘
               │  WebSocket broadcast
  ┌────────────▼──────────────┐
  │  TODOS los navegadores    │
  │  conectados reciben el    │
  │  evento en tiempo real    │
  └───────────────────────────┘"""

box(s, 0.3, 1.15, 8.1, 6.1, bg=SURFACE)
txt(s, arch_text, 0.45, 1.2, 7.8, 6.0,
    size=10.5, color=CYAN, wrap=False)

# Leyenda derecha
leyenda = [
  (PRIMARY, "Frontend HTML/CSS/JS"),
  (SUCCESS, "Base de datos PostgreSQL"),
  (WARNING, "Mensajería Redis"),
  (CYAN,    "Tiempo real Socket.io"),
  (DANGER,  "Autenticación JWT"),
  (YELLOW,  "IaC CloudFormation"),
]
txt(s, "Componentes", 8.7, 1.2, 4.3, 0.4, size=14, bold=True, color=WHITE)
for i, (color, label) in enumerate(leyenda):
    y = 1.75 + i * 0.75
    b = box(s, 8.7, y, 0.35, 0.35, bg=color)
    b.line.fill.background()
    txt(s, label, 9.15, y, 3.8, 0.38, size=12, color=WHITE)

footer(s, 4, 18)

# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 5 — PASO 1: ABRIR SWAGGER
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header_bar(s, "Demo Paso 1 — Swagger UI en Producción",
           "studysync-api-a5o9.onrender.com/api-docs", "API Docs", CYAN)

num_badge(s, "1", 0.35, 1.3)
txt(s, "Abrir en el navegador:", 0.9, 1.28, 8, 0.4, size=15, bold=True, color=WHITE)

box(s, 0.9, 1.72, 10.5, 0.55, bg=SURFACE)
txt(s, "  https://studysync-api-a5o9.onrender.com/api-docs",
    0.9, 1.72, 10.5, 0.55, size=14, color=CYAN, bold=True)

txt(s, "Qué deben ver:", 0.35, 2.4, 5, 0.4, size=13, bold=True, color=MUTED)
items_swagger = [
  "Sección Auth: POST /auth/register y POST /auth/login",
  "Sección Sesiones: 5 endpoints con candado 🔒 en POST/PUT/DELETE",
  "Botón Authorize (candado) en la esquina superior derecha",
  "Servidor: https://studysync-api-a5o9.onrender.com",
]
for i, item in enumerate(items_swagger):
    bullet_row(s, "✓", item, 0.35, 2.85 + i*0.4, size=13)

screenshot_placeholder(s, 6.5, 2.35, 6.5, 4.7,
    "[ Captura: Swagger UI cargado en producción\n  con secciones Auth y Sesiones visibles ]")

footer(s, 5, 18)

# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 6 — PASO 2: REGISTRO
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header_bar(s, "Demo Paso 2 — Registrar Usuarios",
           "POST /auth/register — contraseña hasheada con bcrypt", "bcrypt 12 rounds", WARNING)

num_badge(s, "2", 0.35, 1.3)
txt(s, "Registrar los 3 grupos de estudio:", 0.9, 1.28, 8, 0.4,
    size=15, bold=True, color=WHITE)

grupos = [
  ("Grupo 1", "Ana Torres",     "ana.torres@upds.edu.bo",     "grupo1_2026", "Matemáticas",          WARNING),
  ("Grupo 2", "Carlos Méndez",  "carlos.mendez@upds.edu.bo",  "grupo2_2026", "Redes de Computadoras", CYAN),
  ("Grupo 3", "Sofía Vargas",   "sofia.vargas@upds.edu.bo",   "grupo3_2026", "Programación IV",      SUCCESS),
]

for i, (grp, nombre, email, pwd, materia, color) in enumerate(grupos):
    y = 1.85 + i * 1.4
    box(s, 0.3, y, 5.8, 1.25, bg=SURFACE)
    b = box(s, 0.3, y, 0.12, 1.25, bg=color)
    b.line.fill.background()
    txt(s, grp, 0.55, y+0.05, 1.2, 0.35, size=11, bold=True, color=color)
    txt(s, nombre, 0.55, y+0.38, 5.2, 0.35, size=14, bold=True, color=WHITE)
    txt(s, email,  0.55, y+0.7,  5.2, 0.3,  size=11, color=MUTED)
    txt(s, materia,0.55, y+0.95, 5.2, 0.28, size=11, color=color)

# Respuesta esperada
box(s, 6.4, 1.85, 6.6, 2.5, bg=SURFACE)
txt(s, "Respuesta esperada (201 Created):", 6.5, 1.95, 6.2, 0.35,
    size=12, bold=True, color=WHITE)
codigo_resp = '{\n  "ok": true,\n  "mensaje": "Usuario registrado exitosamente",\n  "usuario": {\n    "id": 3,\n    "nombre": "Carlos Méndez",\n    "email": "carlos.mendez@...",\n    "creadoEn": "2026-..."\n  }\n}'
txt(s, codigo_resp, 6.5, 2.35, 6.3, 1.85, size=10.5, color=CYAN, wrap=False)

txt(s, "⚠️  El campo password NUNCA aparece en la respuesta", 6.4, 4.5, 6.5, 0.4,
    size=12, bold=True, color=WARNING)

screenshot_placeholder(s, 6.4, 4.95, 6.6, 2.2,
    "[ Captura: POST /auth/register\n  con respuesta 201 en Swagger ]")

footer(s, 6, 18)

# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 7 — PASO 3: LOGIN + TOKEN
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header_bar(s, "Demo Paso 3 — Login y Token JWT",
           "POST /auth/login → devuelve eyJhbGci... (JSON Web Token)", "JWT", DANGER)

num_badge(s, "3", 0.35, 1.3)
txt(s, "¿Qué es un JWT?", 0.9, 1.28, 6, 0.4, size=15, bold=True, color=WHITE)

partes = [
  ("HEADER",    "eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9",  "Algoritmo: HS256",        WARNING),
  ("PAYLOAD",   "eyJpZCI6MSwiZW1haWwiOiIuLi4ifQ",         "Datos: id, email, nombre", CYAN),
  ("SIGNATURE", "SflKxwRJSMeKKF2QT4fwpMeJf36POk6yJV_",   "Firma con JWT_SECRET",     SUCCESS),
]
txt(s, ".", 0.35, 1.75, 12.5, 0.3, size=8, color=MUTED)
for i, (parte, valor, desc, color) in enumerate(partes):
    y = 1.9 + i * 0.85
    box(s, 0.35, y, 12.3, 0.75, bg=SURFACE)
    txt(s, parte, 0.5, y+0.08, 1.6, 0.3, size=10, bold=True, color=color)
    txt(s, valor, 2.2, y+0.05, 7.5, 0.32, size=9.5, color=WHITE, wrap=False)
    txt(s, "→ " + desc, 2.2, y+0.38, 7, 0.28, size=10, color=MUTED)

txt(s, "Analogía: es como un ticket de cine firmado — no se puede falsificar,\npero todos pueden leer qué película dice (no guardes passwords en el payload).",
    0.35, 4.55, 8.5, 0.8, size=12, color=YELLOW, italic=True)

txt(s, "Pasos en Swagger:", 0.35, 5.45, 4, 0.35, size=13, bold=True, color=WHITE)
for i, step in enumerate([
    "Expandir POST /auth/login → Try it out",
    "Poner email y password → Execute",
    "Copiar el token de la respuesta",
    "Clic en Authorize (candado) → pegar → Authorize",
]):
    num_badge(s, i+1, 0.35, 5.88 + i*0.38, PRIMARY)
    txt(s, step, 0.9, 5.88 + i*0.38, 5.5, 0.36, size=12, color=WHITE)

screenshot_placeholder(s, 6.5, 4.3, 6.5, 2.9,
    "[ Captura: respuesta del login\n  con el token JWT visible\n  + botón Authorize completado ]")

footer(s, 7, 18)

# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 8 — PASO 4: CREAR SESIONES (GRUPO 1)
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header_bar(s, "Demo Paso 4 — Grupo 1: Ana Torres (Matemáticas)",
           "POST /api/sesiones con Authorization: Bearer <token>", "3 sesiones", WARNING)

num_badge(s, "4", 0.35, 1.3)
txt(s, "Crear las 3 sesiones de Matemáticas:", 0.9, 1.28, 8, 0.4,
    size=15, bold=True, color=WHITE)

sesiones_g1 = [
  ("#11", "Algebra Lineal — Matrices",      "Operaciones con matrices y determinantes"),
  ("#12", "Algebra Lineal — Vectores",      "Espacios vectoriales y transformaciones lineales"),
  ("#13", "Algebra Lineal — Examen parcial","Repaso completo antes del parcial"),
]
for i, (sid, titulo, desc) in enumerate(sesiones_g1):
    y = 1.85 + i * 1.0
    box(s, 0.35, y, 6.0, 0.88, bg=SURFACE)
    b = box(s, 0.35, y, 0.1, 0.88, bg=WARNING)
    b.line.fill.background()
    txt(s, sid, 0.55, y+0.08, 0.7, 0.3, size=11, bold=True, color=WARNING)
    txt(s, titulo, 1.25, y+0.05, 5.0, 0.35, size=13, bold=True, color=WHITE)
    txt(s, desc,   1.25, y+0.45, 5.0, 0.32, size=11, color=MUTED)

txt(s, "Body JSON para cada sesión:", 0.35, 4.98, 5.5, 0.35,
    size=12, bold=True, color=WHITE)
box(s, 0.35, 5.38, 5.9, 1.75, bg=SURFACE)
txt(s, '{\n  "titulo": "Algebra Lineal — Matrices",\n  "materia": "Matematicas",\n  "descripcion": "Operaciones con matrices"\n}',
    0.5, 5.42, 5.6, 1.65, size=11, color=CYAN, wrap=False)

screenshot_placeholder(s, 6.5, 1.85, 6.5, 5.3,
    "[ Captura: POST /api/sesiones con token\n  y respuesta 201 mostrando\n  id, titulo, materia y autor ]")

footer(s, 8, 18)

# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 9 — PASO 5: GRUPO 2 Y 3
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header_bar(s, "Demo Paso 5 — Grupos 2 y 3: Redes y Prog IV",
           "Cada grupo con su propio token JWT", "6 sesiones más", CYAN)

num_badge(s, "5", 0.35, 1.3)

# Grupo 2
box(s, 0.35, 1.3, 6.0, 0.45, bg=RGBColor(0x06,0x40,0x5C))
txt(s, "  Carlos Méndez — Redes de Computadoras", 0.35, 1.3, 6.0, 0.45,
    size=13, bold=True, color=CYAN)
redes = [
  ("#5", "Redes — Modelo OSI",                  "Las 7 capas del modelo OSI explicadas"),
  ("#6", "Redes — TCP/IP y subnetting",          "Cálculo de subredes IPv4"),
  ("#7", "Redes — Laboratorio Cisco Packet Tracer","Configuración de routers y switches"),
]
for i, (sid, titulo, desc) in enumerate(redes):
    y = 1.82 + i * 0.72
    box(s, 0.35, y, 6.0, 0.65, bg=SURFACE)
    b = box(s, 0.35, y, 0.1, 0.65, bg=CYAN)
    b.line.fill.background()
    txt(s, sid, 0.52, y+0.06, 0.6, 0.28, size=10, bold=True, color=CYAN)
    txt(s, titulo, 1.18, y+0.04, 5.0, 0.28, size=12, bold=True, color=WHITE)
    txt(s, desc,   1.18, y+0.35, 5.0, 0.25, size=10, color=MUTED)

# Grupo 3
box(s, 0.35, 4.05, 6.0, 0.45, bg=RGBColor(0x05,0x40,0x20))
txt(s, "  Sofía Vargas — Programación IV", 0.35, 4.05, 6.0, 0.45,
    size=13, bold=True, color=SUCCESS)
prog = [
  ("#8",  "Prog IV — Express REST API",      "Construcción de APIs con Node.js"),
  ("#9",  "Prog IV — Prisma ORM y Supabase", "Conexión a PostgreSQL con Prisma"),
  ("#10", "Prog IV — Redis PubSub y Socket.io","Mensajería en tiempo real"),
]
for i, (sid, titulo, desc) in enumerate(prog):
    y = 4.57 + i * 0.72
    box(s, 0.35, y, 6.0, 0.65, bg=SURFACE)
    b = box(s, 0.35, y, 0.1, 0.65, bg=SUCCESS)
    b.line.fill.background()
    txt(s, sid, 0.52, y+0.06, 0.6, 0.28, size=10, bold=True, color=SUCCESS)
    txt(s, titulo, 1.18, y+0.04, 5.0, 0.28, size=12, bold=True, color=WHITE)
    txt(s, desc,   1.18, y+0.35, 5.0, 0.25, size=10, color=MUTED)

screenshot_placeholder(s, 6.5, 1.3, 6.5, 5.85,
    "[ Captura: GET /api/sesiones mostrando\n  las 12 sesiones de los 3 grupos\n  con autor incluido en cada una ]")

footer(s, 9, 18)

# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 10 — PASO 6: SUPABASE TABLE EDITOR
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header_bar(s, "Demo Paso 6 — Datos Reales en Supabase",
           "Table Editor → tabla 'sesiones' — persistencia real en la nube", "PostgreSQL ☁", SUCCESS)

num_badge(s, "6", 0.35, 1.3)
txt(s, "Verificar en Supabase que los datos persisten:", 0.9, 1.28, 8, 0.4,
    size=15, bold=True, color=WHITE)

txt(s, "Instrucciones:", 0.35, 1.8, 4, 0.35, size=13, bold=True, color=WHITE)
pasos_supa = [
  "Abrir supabase.com → tu proyecto",
  "Table Editor → tabla sesiones",
  "Ver las 12 filas creadas desde la API",
  "Ver tabla usuarios → 5 usuarios registrados",
  "Verificar que 'autor_id' relaciona correctamente",
]
for i, paso in enumerate(pasos_supa):
    num_badge(s, i+1, 0.35, 2.25 + i*0.55, SUCCESS)
    txt(s, paso, 0.9, 2.25 + i*0.55, 5.5, 0.45, size=13, color=WHITE)

txt(s, "Punto pedagógico clave:", 0.35, 5.15, 5.5, 0.35, size=13, bold=True, color=YELLOW)
txt(s, "Antes usábamos un array en memoria [ ].\nSi el servidor se reinicia, se pierden todos los datos.\nAhora con Prisma + Supabase, los datos sobreviven\ncualquier reinicio. Esto es persistencia real.",
    0.35, 5.55, 5.8, 1.2, size=12, color=WHITE)

screenshot_placeholder(s, 6.5, 1.3, 6.5, 3.1,
    "[ Captura: Supabase Table Editor\n  mostrando tabla 'sesiones'\n  con las 12 filas visibles ]")
screenshot_placeholder(s, 6.5, 4.5, 6.5, 2.65,
    "[ Captura: tabla 'usuarios'\n  con los 5 registros\n  (sin el campo password) ]")

footer(s, 10, 18)

# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 11 — PASO 7: TIEMPO REAL (SOCKET.IO)
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header_bar(s, "Demo Paso 7 — Tiempo Real con Redis + Socket.io",
           "El evento viaja: API → Redis → Socket.io → TODOS los navegadores", "Pub/Sub Live", PRIMARY)

num_badge(s, "7", 0.35, 1.3)
txt(s, "Abrir DOS pestañas del navegador:", 0.9, 1.28, 8, 0.4,
    size=15, bold=True, color=WHITE)

# Flujo visual
flujo = [
  ("1", "Pestaña A", "studysync-api-a5o9.onrender.com → login → panel", PRIMARY),
  ("2", "Pestaña B", "studysync-api-a5o9.onrender.com → login con otro usuario", CYAN),
  ("3", "En Pestaña A", "Crear una sesión nueva desde el formulario", WARNING),
  ("4", "En Pestaña B", "El feed de eventos se actualiza INSTANTÁNEAMENTE", SUCCESS),
  ("5", "En ambas",  "Badge '🟢 Conectado' visible — WebSocket activo", SUCCESS),
]
for i, (n, label, desc, color) in enumerate(flujo):
    y = 1.85 + i * 0.88
    box(s, 0.35, y, 6.0, 0.78, bg=SURFACE)
    b = box(s, 0.35, y, 0.1, 0.78, bg=color)
    b.line.fill.background()
    num_badge(s, n, 0.5, y+0.16, color)
    txt(s, label, 1.1, y+0.06, 2.5, 0.3, size=12, bold=True, color=color)
    txt(s, desc,  1.1, y+0.4, 4.9, 0.32, size=11, color=WHITE)

txt(s, "Analogía Redis Pub/Sub:", 0.35, 6.42, 6.0, 0.3, size=12, bold=True, color=YELLOW)
txt(s, '"Como un grupo de WhatsApp: publicas un mensaje y\ntodos los miembros del grupo lo reciben al instante."',
    0.35, 6.72, 6.0, 0.5, size=11, color=YELLOW, italic=True)

screenshot_placeholder(s, 6.5, 1.3, 6.5, 2.9,
    "[ Captura: Pestaña A — formulario\n  y botón crear sesión ]")
screenshot_placeholder(s, 6.5, 4.3, 6.5, 2.9,
    "[ Captura: Pestaña B — feed\n  mostrando el evento recibido\n  en tiempo real ]")

footer(s, 11, 18)

# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 12 — PASO 8: SEGURIDAD (HELMET + JWT)
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header_bar(s, "Demo Paso 8 — Headers de Seguridad (Helmet) + JWT",
           "DevTools F12 → Network → Response Headers", "Helmet 🛡", DANGER)

num_badge(s, "8", 0.35, 1.3)
txt(s, "Verificar con DevTools del navegador (F12 → Network):", 0.9, 1.28, 8, 0.4,
    size=15, bold=True, color=WHITE)

headers_helmet = [
  ("x-frame-options",             "SAMEORIGIN",    "Evita clickjacking (no embed en iframe externo)"),
  ("x-content-type-options",      "nosniff",       "El browser NO adivina el tipo MIME del archivo"),
  ("cross-origin-opener-policy",  "same-origin",   "Aísla el contexto de navegación (anti-Spectre)"),
  ("referrer-policy",             "no-referrer",   "No envía la URL actual como Referer a otros sitios"),
  ("x-dns-prefetch-control",      "off",           "Desactiva la resolución DNS anticipada"),
  ("content-security-policy",     "(configurado)", "Controla qué scripts/estilos puede cargar la página"),
]
txt(s, "Header", 0.35, 1.82, 3.2, 0.35, size=11, bold=True, color=MUTED)
txt(s, "Valor", 3.6, 1.82, 2.5, 0.35, size=11, bold=True, color=MUTED)
txt(s, "Protección", 6.15, 1.82, 6.8, 0.35, size=11, bold=True, color=MUTED)

for i, (h, v, desc) in enumerate(headers_helmet):
    y = 2.22 + i * 0.62
    box(s, 0.35, y, 12.6, 0.55, bg=SURFACE if i%2==0 else BG)
    txt(s, h, 0.5, y+0.1, 3.0, 0.35, size=10.5, color=DANGER, wrap=False)
    txt(s, v, 3.6, y+0.1, 2.4, 0.35, size=10.5, color=WARNING, wrap=False)
    txt(s, desc, 6.15, y+0.1, 6.6, 0.35, size=10.5, color=WHITE)

txt(s, "POST sin token → 401 | Token inválido → 401 | Token expirado → 401",
    0.35, 6.08, 8, 0.38, size=12, color=DANGER, bold=True)

screenshot_placeholder(s, 8.8, 5.7, 4.2, 1.5,
    "[ Captura: DevTools → Headers\n  de seguridad en respuesta ]")

footer(s, 12, 18)

# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 13 — PASO 9: IAC LOCALSTACK
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header_bar(s, "Demo Paso 9 — Infrastructure as Code (IaC)",
           "AWS CloudFormation + LocalStack — mismo YAML, diferente comando", "IaC ☁", YELLOW)

num_badge(s, "9", 0.35, 1.3)
txt(s, "¿Qué es IaC? La infraestructura como código:", 0.9, 1.28, 8, 0.4,
    size=15, bold=True, color=WHITE)

txt(s, "Sin IaC (clickops):", 0.35, 1.82, 4, 0.3, size=12, bold=True, color=DANGER)
for i, item in enumerate(["Entrar a la consola AWS", "Clic clic clic...", "Diferente cada vez", "No se puede reproducir"]):
    bullet_row(s, "✗", item, 0.35, 2.18 + i*0.38, size=12, icon_color=DANGER)

txt(s, "Con IaC (CloudFormation):", 4.0, 1.82, 4, 0.3, size=12, bold=True, color=SUCCESS)
for i, item in enumerate(["Escribir template.yaml", "Un comando lo crea todo", "Idéntico en dev/staging/prod", "Versionado en Git"]):
    bullet_row(s, "✓", item, 4.0, 2.18 + i*0.38, size=12, icon_color=SUCCESS)

txt(s, "Recursos creados en LocalStack:", 0.35, 4.1, 6, 0.35, size=13, bold=True, color=WHITE)
recursos = [
  ("S3 Bucket",        "studysync-adjuntos-development",  "Almacenamiento de archivos adjuntos", WARNING),
  ("DynamoDB Table",   "studysync-eventos-development",   "Registro de eventos con TTL",         CYAN),
  ("SSM Parameter",    "/studysync/development/api-url",  "Configuración centralizada",           SUCCESS),
]
for i, (tipo, nombre, desc, color) in enumerate(recursos):
    y = 4.52 + i * 0.78
    box(s, 0.35, y, 6.0, 0.68, bg=SURFACE)
    b = box(s, 0.35, y, 0.1, 0.68, bg=color)
    b.line.fill.background()
    txt(s, tipo, 0.55, y+0.05, 2.2, 0.28, size=11, bold=True, color=color)
    txt(s, nombre, 2.8, y+0.05, 3.4, 0.28, size=10.5, color=WHITE, wrap=False)
    txt(s, desc, 0.55, y+0.38, 5.5, 0.25, size=10, color=MUTED)

screenshot_placeholder(s, 6.5, 4.05, 6.5, 3.1,
    "[ Captura: terminal con resultado\n  awslocal cloudformation\n  describe-stack-resources --output table ]")

box(s, 0.35, 2.15, 0.04, 1.5, bg=MUTED)

footer(s, 13, 18)

# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 14 — PASO 10: REPOSITORIO GITHUB
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header_bar(s, "Demo Paso 10 — Repositorio GitHub",
           "github.com/bolivianotech/studysync-api", "Open Source", PRIMARY)

num_badge(s, "10", 0.35, 1.3)
txt(s, "Lo que deben ver en el repositorio:", 0.9, 1.28, 8, 0.4,
    size=15, bold=True, color=WHITE)

commits = [
  ("0c14d07", "docs: README completo y .env.example para nuevos alumnos"),
  ("6e4b64a", "chore: ignorar directorio .claude en git"),
  ("09fe534", "docs: guía final IaC, Prisma+Supabase y producción"),
  ("c79fecc", "chore: agregar evento de prueba DynamoDB para demo"),
  ("d0e271f", "feat: JWT auth, Prisma+Supabase, Helmet, frontend, CloudFormation"),
  ("f7fca67", "fix: corregir puerto Render y completar Redis pub/sub"),
  ("8516f99", "feat: swagger UI, CORS, rate-limit, redis pub/sub, socket.io"),
  ("945da3a", "feat: CRUD completo de sesiones con 5 endpoints"),
  ("01413ee", "feat: setup inicial del proyecto StudySync"),
]
txt(s, "Historial de commits (10 commits):", 0.35, 1.82, 6, 0.35, size=12, bold=True, color=WHITE)
for i, (sha, msg) in enumerate(commits):
    y = 2.22 + i * 0.5
    txt(s, sha, 0.4, y, 1.1, 0.38, size=9.5, color=WARNING, wrap=False)
    txt(s, msg, 1.6, y, 4.8, 0.38, size=9.5, color=WHITE)

screenshot_placeholder(s, 6.5, 1.3, 6.5, 3.3,
    "[ Captura: GitHub — árbol de archivos\n  con estructura src/, prisma/,\n  cloudformation/, public/ ]")
screenshot_placeholder(s, 6.5, 4.75, 6.5, 2.4,
    "[ Captura: GitHub — historial\n  de commits con mensajes\n  descriptivos ]")

footer(s, 14, 18)

# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 15 — CHECKLIST QA
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header_bar(s, "Checklist de Verificación Final QA",
           "16 puntos — todos deben estar en verde para aprobar", "16/16 ✓", SUCCESS)

checks = [
  ("/register crea usuario → id + sin password",        True),
  ("/register email duplicado → 409 Conflict",          True),
  ("/login correcto → JWT en respuesta",                 True),
  ("/login password incorrecto → 401 genérico",         True),
  ("POST sin token → 401 con hint",                     True),
  ("DELETE sin token → 401",                            True),
  ("GET /api/sesiones sin token → 200 público",         True),
  ("POST con JWT → 201, sesión en Supabase",            True),
  ("PUT con JWT → campos actualizados en Supabase",     True),
  ("DELETE con JWT → eliminado de Supabase",            True),
  ("GET /999 → 404 not found",                          True),
  ("Swagger carga con BearerAuth y 4 paths",            True),
  ("Frontend: badge 🟢 Conectado visible",             True),
  ("Crear sesión → feed actualizado en 2 pestañas",    True),
  ("Headers Helmet en respuestas (X-Frame-Options...)", True),
  ("Deploy Render → ✓ Deployed, sin PORT fijo",        True),
]

col_size = 8
for i, (item, ok) in enumerate(checks):
    col = i % 2
    row = i // 2
    x = 0.3 + col * 6.5
    y = 1.28 + row * 0.72
    box(s, x, y, 6.2, 0.62, bg=SURFACE)
    icon = "✓" if ok else "✗"
    color = SUCCESS if ok else DANGER
    txt(s, icon, x+0.12, y+0.1, 0.4, 0.38, size=16, bold=True, color=color)
    txt(s, item, x+0.6, y+0.13, 5.4, 0.38, size=11, color=WHITE)

footer(s, 15, 18)

# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 16 — COBERTURA DE ACTIVIDADES
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header_bar(s, "Cobertura — Actividades 3 y 4",
           "Todo lo que aprendieron en este proyecto", "Aprobado ✓")

txt(s, "ACTIVIDAD 3 — API REST Segura", 0.35, 1.22, 6.3, 0.45,
    size=17, bold=True, color=PRIMARY)
box(s, 0.35, 1.22, 0.08, 5.3, bg=PRIMARY)

act3 = [
  ("API REST completa",          "5 endpoints CRUD + 2 de autenticación"),
  ("Base de datos real",         "Supabase PostgreSQL con Prisma ORM"),
  ("Autenticación JWT",          "Register, Login, middleware autenticar.js"),
  ("Seguridad HTTP",             "Helmet: X-Frame-Options, CSP, HSTS y más"),
  ("CORS + Rate Limiting",       "100 req/15min por IP en /api/"),
  ("Documentación OpenAPI 3.0",  "Swagger UI con BearerAuth funcional"),
]
for i, (titulo, desc) in enumerate(act3):
    y = 1.75 + i * 0.72
    txt(s, "✓  " + titulo, 0.55, y, 3.0, 0.32, size=12, bold=True, color=SUCCESS)
    txt(s, desc, 0.55, y+0.32, 5.7, 0.28, size=11, color=MUTED)

txt(s, "ACTIVIDAD 4 — Tiempo Real e IaC", 6.7, 1.22, 6.3, 0.45,
    size=17, bold=True, color=CYAN)
box(s, 6.7, 1.22, 0.08, 5.3, bg=CYAN)

act4 = [
  ("Redis Pub/Sub",              "Upstash → eventos publicados en cada creación"),
  ("WebSockets Socket.io",       "Broadcast a todos los clientes conectados"),
  ("Frontend en tiempo real",    "Badge conexión + feed de eventos en vivo"),
  ("Infrastructure as Code",     "CloudFormation YAML con S3, DynamoDB, SSM"),
  ("LocalStack",                 "Simula AWS en Docker — mismo YAML, sin costo"),
  ("Deploy CI/CD",               "GitHub push → Render despliega automáticamente"),
]
for i, (titulo, desc) in enumerate(act4):
    y = 1.75 + i * 0.72
    txt(s, "✓  " + titulo, 6.9, y, 3.2, 0.32, size=12, bold=True, color=SUCCESS)
    txt(s, desc, 6.9, y+0.32, 5.9, 0.28, size=11, color=MUTED)

footer(s, 16, 18)

# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 17 — PRÓXIMOS PASOS
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header_bar(s, "Próximos Pasos — ¿Qué pueden hacer el fin de semana?",
           "El repositorio es tuyo — clónalo, modifícalo, defiéndelo", "Weekend Task 🚀")

txt(s, "Para el fin de semana (defensa el Lunes):", 0.35, 1.25, 9, 0.4,
    size=16, bold=True, color=WHITE)

tareas = [
  ("Clonar el repositorio",
   "git clone https://github.com/bolivianotech/studysync-api.git",
   "Tener el código localmente", PRIMARY),
  ("Crear tus propias credenciales",
   "Supabase gratis (2 proyectos) + Upstash gratis (10k cmd/día)",
   "Tu propia instancia del proyecto", CYAN),
  ("Hacer npm install y npm run dev",
   "El servidor local debe levantar en http://localhost:3000",
   "Verificar instalación local", SUCCESS),
  ("Crear sesiones con TUS datos",
   "Registrar 3 usuarios distintos, crear 9+ sesiones reales",
   "Demostrar que entiendes el flujo", WARNING),
  ("Modificar algo del proyecto",
   "Agregar un campo, cambiar el CSS, añadir un endpoint nuevo",
   "Demostrar apropiación del código", YELLOW),
  ("Preparar la demo de 5 minutos",
   "Swagger + Supabase + Frontend en tiempo real + GitHub",
   "Practicar el guión de presentación", DANGER),
]

for i, (accion, comando, objetivo, color) in enumerate(tareas):
    col = i % 2
    row = i // 2
    x = 0.3 + col * 6.5
    y = 1.82 + row * 1.68
    box(s, x, y, 6.2, 1.55, bg=SURFACE)
    b = box(s, x, y, 0.1, 1.55, bg=color)
    b.line.fill.background()
    txt(s, accion, x+0.25, y+0.08, 5.7, 0.35, size=13, bold=True, color=WHITE)
    box(s, x+0.25, y+0.48, 5.7, 0.52, bg=BG)
    txt(s, comando, x+0.35, y+0.52, 5.5, 0.42, size=9.5, color=color, wrap=False)
    txt(s, "→ " + objetivo, x+0.25, y+1.12, 5.7, 0.32, size=10.5, color=MUTED)

footer(s, 17, 18)

# ══════════════════════════════════════════════════════════════════════════════
# DIAPOSITIVA 18 — CIERRE
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
box(s, 0, 0, 13.33, 7.5, bg=SURFACE)
box(s, 0, 3.68, 13.33, 0.08, bg=PRIMARY)

txt(s, "¡Proyecto Completado!", 0.5, 0.35, 12.3, 1.1,
    size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "StudySync API — en producción, en la nube, con datos reales",
    0.5, 1.45, 12.3, 0.6, size=18, color=MUTED, align=PP_ALIGN.CENTER)

txt(s, "📦  github.com/bolivianotech/studysync-api",
    0.5, 2.22, 12.3, 0.55, size=17, color=PRIMARY, align=PP_ALIGN.CENTER, bold=True)
txt(s, "🌐  studysync-api-a5o9.onrender.com",
    0.5, 2.78, 12.3, 0.55, size=17, color=CYAN, align=PP_ALIGN.CENTER, bold=True)
txt(s, "📚  /api-docs  (Swagger UI con todos los endpoints)",
    0.5, 3.22, 12.3, 0.45, size=14, color=SUCCESS, align=PP_ALIGN.CENTER)

txt(s, "Lo que construyeron:", 0.5, 3.9, 12.3, 0.4,
    size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

logros = "API REST  •  JWT Auth  •  Prisma ORM  •  Supabase PostgreSQL  •  Redis Pub/Sub\nSocket.io WebSockets  •  Helmet Security  •  Swagger OpenAPI  •  CloudFormation IaC\nLocalStack Docker  •  Render Deploy  •  Frontend Vanilla JS  •  Git CI/CD"
txt(s, logros, 0.5, 4.35, 12.3, 1.2,
    size=13, color=YELLOW, align=PP_ALIGN.CENTER)

txt(s, "M.Sc. Jimmy Nataniel Requena Llorentty  •  Programación IV — UPDS 2026",
    0.5, 6.5, 12.3, 0.45, size=12, color=MUTED, align=PP_ALIGN.CENTER)

# ── Guardar ───────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Presentacion generada: {OUT}")
print(f"Diapositivas: {len(prs.slides)}")
